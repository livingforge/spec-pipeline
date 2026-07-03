"""共有フィクスチャとテスト用ドキュメントのビルダー。

各形式 (docx/xlsx/pptx/pdf) を実ファイルとしてプログラム的に生成し、
抽出器を end-to-end で検証できるようにする。
"""

from __future__ import annotations

from pathlib import Path

import fitz  # PyMuPDF
import pytest


@pytest.fixture(scope="session")
def png_bytes() -> bytes:
    """小さな有効な PNG バイト列 (6x6, 赤)。"""
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 6, 6))
    pix.set_rect(pix.irect, (200, 50, 50))
    return pix.tobytes("png")


@pytest.fixture
def png_file(tmp_path: Path, png_bytes: bytes) -> Path:
    p = tmp_path / "fixture.png"
    p.write_bytes(png_bytes)
    return p


# --------------------------------------------------------------------------
# docx ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_docx(tmp_path: Path):
    from docx import Document

    def _make(
        name: str = "doc.docx",
        *,
        paragraphs: list[tuple[str, str | None]] | None = None,
        table: list[list[str]] | None = None,
        image_path: Path | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        d = Document()
        if title is not None:
            d.core_properties.title = title
        if author is not None:
            d.core_properties.author = author
        for text, style in paragraphs or []:
            if style:
                d.add_paragraph(text, style=style)
            else:
                d.add_paragraph(text)
        if table is not None:
            n_rows = len(table)
            n_cols = max((len(r) for r in table), default=0)
            t = d.add_table(rows=n_rows, cols=n_cols)
            for i, row in enumerate(table):
                for j, val in enumerate(row):
                    t.rows[i].cells[j].text = val
        if image_path is not None:
            d.add_picture(str(image_path))
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        d.save(str(path))
        return path

    return _make


# --------------------------------------------------------------------------
# xlsx ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_xlsx(tmp_path: Path):
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    def _make(
        name: str = "book.xlsx",
        *,
        sheets: dict[str, list[list]] | None = None,
        merges: dict[str, list[str]] | None = None,  # sheet -> ["B1:B3", ...]
        image: tuple[str, Path, str] | None = None,  # (sheet, png_path, anchor)
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        wb = Workbook()
        # デフォルトシートを消す前に少なくとも 1 枚必要
        default = wb.active
        created_any = False
        for sheet_name, grid in (sheets or {}).items():
            if not created_any:
                ws = default
                ws.title = sheet_name
                created_any = True
            else:
                ws = wb.create_sheet(sheet_name)
            for r, row in enumerate(grid, start=1):
                for c, val in enumerate(row, start=1):
                    ws.cell(row=r, column=c, value=val)
        if not created_any:
            default.title = "Sheet1"
        for sheet_name, refs in (merges or {}).items():
            for ref in refs:
                wb[sheet_name].merge_cells(ref)
        if title is not None:
            wb.properties.title = title
        if author is not None:
            wb.properties.creator = author
        if image is not None:
            sheet_name, png_path, anchor = image
            ws = wb[sheet_name]
            ws.add_image(XLImage(str(png_path)), anchor)
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(path))
        return path

    return _make


# --------------------------------------------------------------------------
# pptx ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_pptx(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    def _make(
        name: str = "deck.pptx",
        *,
        slides: list[dict] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for spec in slides or []:
            slide = prs.slides.add_slide(blank)
            for text in spec.get("texts", []):
                tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
                tb.text_frame.text = text
            for tbl in spec.get("tables", []):
                n_rows = len(tbl)
                n_cols = max((len(r) for r in tbl), default=0)
                gt = slide.shapes.add_table(
                    n_rows, n_cols, Inches(1), Inches(3), Inches(5), Inches(2)
                ).table
                for i, row in enumerate(tbl):
                    for j, val in enumerate(row):
                        gt.cell(i, j).text = val
            for img in spec.get("images", []):
                slide.shapes.add_picture(str(img), Inches(5), Inches(1))
            notes = spec.get("notes")
            if notes is not None:
                slide.notes_slide.notes_text_frame.text = notes
        if title is not None:
            prs.core_properties.title = title
        if author is not None:
            prs.core_properties.author = author
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        return path

    return _make


# --------------------------------------------------------------------------
# pdf ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_pdf(tmp_path: Path):
    def _make(
        name: str = "doc.pdf",
        *,
        pages: list[dict] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        doc = fitz.open()
        for spec in pages or [{}]:
            pg = doc.new_page()
            for text, pos in spec.get("texts", []):
                pg.insert_text(pos, text, fontsize=11)
            # 罫線付きの表を描く: grid = 行数・列数・原点・セルサイズ・セルテキスト
            grid = spec.get("grid")
            if grid:
                n_rows = grid["rows"]
                n_cols = grid["cols"]
                x0, y0 = grid.get("origin", (100, 200))
                cw = grid.get("cw", 80)
                ch = grid.get("ch", 30)
                for i in range(n_rows + 1):
                    pg.draw_line((x0, y0 + i * ch), (x0 + n_cols * cw, y0 + i * ch))
                for j in range(n_cols + 1):
                    pg.draw_line((x0 + j * cw, y0), (x0 + j * cw, y0 + n_rows * ch))
                for (i, j), val in grid.get("cells", {}).items():
                    pg.insert_text((x0 + j * cw + 4, y0 + i * ch + 18), val, fontsize=10)
            for img, rect in spec.get("images", []):
                pg.insert_image(fitz.Rect(*rect), stream=img)
        if title is not None:
            doc.set_metadata({**(doc.metadata or {}), "title": title})
        if author is not None:
            doc.set_metadata({**(doc.metadata or {}), "author": author})
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(path))
        doc.close()
        return path

    return _make
