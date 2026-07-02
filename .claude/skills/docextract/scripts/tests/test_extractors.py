"""形式別抽出器のエンドツーエンドテスト。

フィクスチャ (docx / xlsx / pptx / pdf) はテスト実行時に生成するため、
バイナリファイルの同梱もネットワークも不要。OCR・画像内表検出は
モデルダウンロードを要するため無効化して実行する (純粋な抽出のみ検証)。
"""

import io
import json
import sys
import tempfile
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from docextract import extract


def _make_png() -> bytes:
    from PIL import Image

    img = Image.new("RGB", (40, 20), (200, 40, 40))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _build_minimal_pdf() -> bytes:
    """テキスト 2 ブロックと罫線付きの表を含む最小の PDF を手組みで生成する。"""
    ops = [b"1 w"]
    # 段落ブロック A (行間 15pt → 1 ブロックに結合される)
    ops.append(b"BT /F1 11 Tf 72 720 Td (First line of paragraph.) Tj ET")
    ops.append(b"BT /F1 11 Tf 72 705 Td (Second line continues.) Tj ET")
    # 段落ブロック B (45pt 空ける → 別ブロック)
    ops.append(b"BT /F1 11 Tf 72 660 Td (Standalone second paragraph.) Tj ET")
    # 罫線の表 3 列 x 2 行 (x: 72..372, y: 560..608)
    for y in (608, 584, 560):
        ops.append(b"72 %d m 372 %d l S" % (y, y))
    for x in (72, 172, 272, 372):
        ops.append(b"%d 560 m %d 608 l S" % (x, x))
    cells = [(590, [b"Item", b"Q1", b"Q2"]), (566, [b"Sales", b"100", b"110"])]
    for baseline, row in cells:
        for i, text in enumerate(row):
            ops.append(
                b"BT /F1 10 Tf %d %d Td (%s) Tj ET" % (77 + i * 100, baseline, text)
            )
    content = b"\n".join(ops)

    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length %d >>\nstream\n%s\nendstream" % (len(content), content),
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, body in enumerate(objs, start=1):
        offsets.append(len(out))
        out += b"%d 0 obj\n%s\nendobj\n" % (i, body)
    xref_pos = len(out)
    out += b"xref\n0 %d\n0000000000 65535 f \n" % (len(objs) + 1)
    for off in offsets:
        out += b"%010d 00000 n \n" % off
    out += b"trailer\n<< /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF\n" % (
        len(objs) + 1,
        xref_pos,
    )
    return bytes(out)


class ExtractorTestBase(unittest.TestCase):
    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self.tmp = Path(self._tmp.name)
        self.addCleanup(self._tmp.cleanup)

    def run_extract(self, path: Path) -> dict:
        # OCR と画像内表検出はモデル取得が必要なため無効化 (抽出のみを検証)
        return extract(path, output_dir=self.tmp / "out", ocr=False, image_tables=False)

    def elements_of(self, data: dict, type_: str) -> list[dict]:
        return [e for e in data["elements"] if e["type"] == type_]


class TestDocx(ExtractorTestBase):
    def setUp(self):
        super().setUp()
        from docx import Document
        from docx.shared import Inches

        png = self.tmp / "img.png"
        png.write_bytes(_make_png())
        doc = Document()
        doc.add_heading("見出し", level=1)
        doc.add_paragraph("本文の段落です。")
        t = doc.add_table(rows=2, cols=2)
        for j, v in enumerate(["項目", "値"]):
            t.rows[0].cells[j].text = v
        for j, v in enumerate(["売上", "100"]):
            t.rows[1].cells[j].text = v
        doc.add_picture(str(png), width=Inches(1))
        self.path = self.tmp / "sample.docx"
        doc.save(str(self.path))

    def test_docx_extraction(self):
        data = self.run_extract(self.path)
        self.assertEqual(data["file_type"], "docx")

        texts = self.elements_of(data, "text")
        self.assertEqual([t["content"] for t in texts], ["見出し", "本文の段落です。"])
        self.assertEqual(texts[0]["style"], "Heading 1")

        tables = self.elements_of(data, "table")
        self.assertEqual(tables[0]["rows"], [["項目", "値"], ["売上", "100"]])

        images = self.elements_of(data, "image")
        self.assertEqual(len(images), 1)
        saved = self.tmp / "out" / "sample_docx" / images[0]["file"]
        self.assertTrue(saved.is_file())

        # 出現順: 見出し < 段落 < 表 < 画像
        orders = [e["location"]["order"] for e in data["elements"]]
        self.assertEqual(orders, sorted(orders))

    def test_result_json_written_utf8(self):
        self.run_extract(self.path)
        raw = (self.tmp / "out" / "sample_docx" / "result.json").read_text(
            encoding="utf-8"
        )
        self.assertIn("見出し", raw)  # ensure_ascii=False
        self.assertEqual(json.loads(raw)["source"], "sample.docx")


class TestXlsx(ExtractorTestBase):
    def setUp(self):
        super().setUp()
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image as XlImage

        png = self.tmp / "img.png"
        png.write_bytes(_make_png())
        wb = Workbook()
        ws = wb.active
        ws.title = "売上"
        ws.append(["月", "売上"])
        ws.append(["4月", 100])
        ws.append(["", ""])  # 末尾の空行はトリムされる
        ws.add_image(XlImage(str(png)), "E2")
        self.path = self.tmp / "book.xlsx"
        wb.save(str(self.path))

    def test_xlsx_extraction(self):
        data = self.run_extract(self.path)
        tables = self.elements_of(data, "table")
        self.assertEqual(tables[0]["rows"], [["月", "売上"], ["4月", "100"]])
        self.assertEqual(tables[0]["location"]["sheet"], "売上")

        images = self.elements_of(data, "image")
        self.assertEqual(images[0]["location"]["anchor"], "E2")
        self.assertIn("売上", data["metadata"]["sheets"])


class TestPptx(ExtractorTestBase):
    def setUp(self):
        super().setUp()
        from pptx import Presentation
        from pptx.util import Inches

        png = self.tmp / "img.png"
        png.write_bytes(_make_png())
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # タイトルのみ
        slide.shapes.title.text = "タイトル"
        tbl = slide.shapes.add_table(
            2, 2, Inches(1), Inches(2), Inches(4), Inches(1)
        ).table
        tbl.cell(0, 0).text = "A"
        tbl.cell(0, 1).text = "B"
        tbl.cell(1, 0).text = "1"
        tbl.cell(1, 1).text = "2"
        slide.shapes.add_picture(str(png), Inches(1), Inches(4))
        slide.notes_slide.notes_text_frame.text = "発表者ノート"
        self.path = self.tmp / "deck.pptx"
        prs.save(str(self.path))

    def test_pptx_extraction(self):
        data = self.run_extract(self.path)
        texts = self.elements_of(data, "text")
        contents = [t["content"] for t in texts]
        self.assertIn("タイトル", contents)

        notes = [t for t in texts if t.get("style") == "notes"]
        self.assertEqual(notes[0]["content"], "発表者ノート")

        tables = self.elements_of(data, "table")
        self.assertEqual(tables[0]["rows"], [["A", "B"], ["1", "2"]])
        self.assertEqual(tables[0]["location"]["slide"], 1)

        images = self.elements_of(data, "image")
        self.assertEqual(len(images), 1)
        self.assertEqual(images[0]["width"], 40)


class TestPdf(ExtractorTestBase):
    def setUp(self):
        super().setUp()
        self.path = self.tmp / "doc.pdf"
        self.path.write_bytes(_build_minimal_pdf())

    def test_pdf_extraction(self):
        data = self.run_extract(self.path)
        self.assertEqual(data["metadata"]["page_count"], 1)

        tables = self.elements_of(data, "table")
        self.assertEqual(len(tables), 1)
        self.assertEqual(tables[0]["rows"], [["Item", "Q1", "Q2"], ["Sales", "100", "110"]])
        self.assertEqual(tables[0]["location"]["page"], 1)

        texts = self.elements_of(data, "text")
        self.assertEqual(len(texts), 2)  # 近接 2 行は 1 ブロックに結合
        self.assertEqual(
            texts[0]["content"],
            "First line of paragraph.\nSecond line continues.",
        )
        self.assertEqual(texts[1]["content"], "Standalone second paragraph.")
        # 表のセル文字列はテキスト要素に重複しない
        for t in texts:
            self.assertNotIn("Sales", t["content"])


if __name__ == "__main__":
    unittest.main()
