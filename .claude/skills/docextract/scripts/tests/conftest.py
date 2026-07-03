"""共有フィクスチャとテスト用ドキュメントのビルダー。

各形式 (docx/xlsx/pptx/pdf) を実ファイルとしてプログラム的に生成し、
抽出器を end-to-end で検証できるようにする。
"""

from __future__ import annotations

from pathlib import Path

import fitz  # PyMuPDF
import pytest


@pytest.fixture(scope="session")
def png_bytes() -> bytes:
    """小さな有効な PNG バイト列 (6x6, 赤)。"""
    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, 6, 6))
    pix.set_rect(pix.irect, (200, 50, 50))
    return pix.tobytes("png")


@pytest.fixture
def png_file(tmp_path: Path, png_bytes: bytes) -> Path:
    p = tmp_path / "fixture.png"
    p.write_bytes(png_bytes)
    return p


# --------------------------------------------------------------------------
# docx ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_docx(tmp_path: Path):
    from docx import Document

    def _make(
        name: str = "doc.docx",
        *,
        paragraphs: list[tuple[str, str | None]] | None = None,
        table: list[list[str]] | None = None,
        image_path: Path | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        d = Document()
        if title is not None:
            d.core_properties.title = title
        if author is not None:
            d.core_properties.author = author
        for text, style in paragraphs or []:
            if style:
                d.add_paragraph(text, style=style)
            else:
                d.add_paragraph(text)
        if table is not None:
            n_rows = len(table)
            n_cols = max((len(r) for r in table), default=0)
            t = d.add_table(rows=n_rows, cols=n_cols)
            for i, row in enumerate(table):
                for j, val in enumerate(row):
                    t.rows[i].cells[j].text = val
        if image_path is not None:
            d.add_picture(str(image_path))
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        d.save(str(path))
        return path

    return _make


# --------------------------------------------------------------------------
# xlsx ビルダー
# --------------------------------------------------------------------------
def _inject_shapes(
    path: Path, shapes: dict[str, list[dict]], order: list[str]
) -> None:
    """保存済み xlsx にオートシェイプ/コネクタの drawing を直接注入する。

    openpyxl は autoshape を書き出せないため、zip を展開して
    xl/drawings/drawingK.xml と関連リレーションを手で組み立てる。
    各シェイプ dict: {"name","text","cell":(col,row)} = テキスト図形、
    {"connector": True} = コネクタ (テキスト無し、抽出対象外の確認用)。
    """
    import os
    import shutil
    import zipfile

    XDR = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
    AMAIN = "http://schemas.openxmlformats.org/drawingml/2006/main"
    RNS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    def _sp(idx: int, s: dict) -> str:
        col, row = s.get("cell", (0, 0))
        to_col, to_row = s.get("to_cell", (col + 2, row + 2))
        anchor_from = (
            f"<xdr:from><xdr:col>{col}</xdr:col><xdr:colOff>0</xdr:colOff>"
            f"<xdr:row>{row}</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:from>"
            f"<xdr:to><xdr:col>{to_col}</xdr:col><xdr:colOff>0</xdr:colOff>"
            f"<xdr:row>{to_row}</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:to>"
        )
        if s.get("connector"):
            # st/end に接続先シェイプ id を渡すと明示接続 (<a:stCxn>/<a:endCxn>) になる
            cxn = ""
            if s.get("st") is not None:
                cxn += f'<a:stCxn id="{s["st"]}" idx="0"/>'
            if s.get("end") is not None:
                cxn += f'<a:endCxn id="{s["end"]}" idx="0"/>'
            body = (
                f'<xdr:cxnSp macro=""><xdr:nvCxnSpPr>'
                f'<xdr:cNvPr id="{idx}" name="{s.get("name","Conn")}"/>'
                f"<xdr:cNvCxnSpPr>{cxn}</xdr:cNvCxnSpPr></xdr:nvCxnSpPr>"
                f'<xdr:spPr><a:prstGeom prst="straightConnector1"><a:avLst/>'
                f"</a:prstGeom></xdr:spPr></xdr:cxnSp>"
            )
        else:
            paras = "".join(
                f"<a:p><a:r><a:t>{line}</a:t></a:r></a:p>"
                for line in str(s["text"]).split("\n")
            )
            body = (
                f'<xdr:sp macro="" textlink=""><xdr:nvSpPr>'
                f'<xdr:cNvPr id="{idx}" name="{s.get("name","Shape")}"/>'
                f"<xdr:cNvSpPr/></xdr:nvSpPr>"
                f'<xdr:spPr><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></xdr:spPr>'
                f"<xdr:txBody><a:bodyPr/>{paras}</xdr:txBody></xdr:sp>"
            )
        return f"<xdr:twoCellAnchor>{anchor_from}{body}<xdr:clientData/></xdr:twoCellAnchor>"

    tmpd = path.parent / (path.stem + "_unzip")
    if tmpd.exists():
        shutil.rmtree(tmpd)
    with zipfile.ZipFile(path) as z:
        z.extractall(tmpd)

    ct_path = tmpd / "[Content_Types].xml"
    ct = ct_path.read_text(encoding="utf-8")
    sheet_index = {title: i + 1 for i, title in enumerate(order)}
    k = 0
    for sheet_name, sps in shapes.items():
        k += 1
        n = sheet_index[sheet_name]
        drawing = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<xdr:wsDr xmlns:xdr="{XDR}" xmlns:a="{AMAIN}">'
            + "".join(_sp(i + 2, s) for i, s in enumerate(sps))
            + "</xdr:wsDr>"
        )
        (tmpd / "xl" / "drawings").mkdir(parents=True, exist_ok=True)
        (tmpd / "xl" / "drawings" / f"drawing{k}.xml").write_text(drawing, encoding="utf-8")

        sheet_path = tmpd / "xl" / "worksheets" / f"sheet{n}.xml"
        s = sheet_path.read_text(encoding="utf-8")
        if "xmlns:r=" not in s:
            s = s.replace("<worksheet ", f'<worksheet xmlns:r="{RNS}" ', 1)
        s = s.replace("</worksheet>", '<drawing r:id="rIdDraw"/></worksheet>')
        sheet_path.write_text(s, encoding="utf-8")

        rels_dir = tmpd / "xl" / "worksheets" / "_rels"
        rels_dir.mkdir(parents=True, exist_ok=True)
        rels_path = rels_dir / f"sheet{n}.xml.rels"
        rel = (
            f'<Relationship Id="rIdDraw" Type="{RNS}/drawing" '
            f'Target="../drawings/drawing{k}.xml"/>'
        )
        if rels_path.exists():
            r = rels_path.read_text(encoding="utf-8").replace(
                "</Relationships>", rel + "</Relationships>"
            )
        else:
            r = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/'
                f'package/2006/relationships">{rel}</Relationships>'
            )
        rels_path.write_text(r, encoding="utf-8")

        override = (
            f'<Override PartName="/xl/drawings/drawing{k}.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.drawing+xml"/>'
        )
        if f"/xl/drawings/drawing{k}.xml" not in ct:
            ct = ct.replace("</Types>", override + "</Types>")
    ct_path.write_text(ct, encoding="utf-8")

    path.unlink()
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(tmpd):
            for f in files:
                fp = Path(root) / f
                z.write(fp, fp.relative_to(tmpd).as_posix())
    shutil.rmtree(tmpd)


@pytest.fixture
def make_xlsx(tmp_path: Path):
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    def _make(
        name: str = "book.xlsx",
        *,
        sheets: dict[str, list[list]] | None = None,
        merges: dict[str, list[str]] | None = None,  # sheet -> ["B1:B3", ...]
        image: tuple[str, Path, str] | None = None,  # (sheet, png_path, anchor)
        # sheet -> [{"name","text","cell":(col,row)} | {"connector":True}, ...]
        # openpyxl は autoshape を書けないため drawing XML を直接注入する。
        shapes: dict[str, list[dict]] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        wb = Workbook()
        # デフォルトシートを消す前に少なくとも 1 枚必要
        default = wb.active
        created_any = False
        for sheet_name, grid in (sheets or {}).items():
            if not created_any:
                ws = default
                ws.title = sheet_name
                created_any = True
            else:
                ws = wb.create_sheet(sheet_name)
            for r, row in enumerate(grid, start=1):
                for c, val in enumerate(row, start=1):
                    ws.cell(row=r, column=c, value=val)
        if not created_any:
            default.title = "Sheet1"
        for sheet_name, refs in (merges or {}).items():
            for ref in refs:
                wb[sheet_name].merge_cells(ref)
        if title is not None:
            wb.properties.title = title
        if author is not None:
            wb.properties.creator = author
        if image is not None:
            sheet_name, png_path, anchor = image
            ws = wb[sheet_name]
            ws.add_image(XLImage(str(png_path)), anchor)
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(path))
        if shapes:
            order = wb.sheetnames  # sheetN.xml は作成順 = この順
            _inject_shapes(path, shapes, order)
        return path

    return _make


# --------------------------------------------------------------------------
# pptx ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_pptx(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    def _make(
        name: str = "deck.pptx",
        *,
        slides: list[dict] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for spec in slides or []:
            slide = prs.slides.add_slide(blank)
            for text in spec.get("texts", []):
                tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
                tb.text_frame.text = text
            for tbl in spec.get("tables", []):
                n_rows = len(tbl)
                n_cols = max((len(r) for r in tbl), default=0)
                gt = slide.shapes.add_table(
                    n_rows, n_cols, Inches(1), Inches(3), Inches(5), Inches(2)
                ).table
                for i, row in enumerate(tbl):
                    for j, val in enumerate(row):
                        gt.cell(i, j).text = val
            for img in spec.get("images", []):
                slide.shapes.add_picture(str(img), Inches(5), Inches(1))
            notes = spec.get("notes")
            if notes is not None:
                slide.notes_slide.notes_text_frame.text = notes
        if title is not None:
            prs.core_properties.title = title
        if author is not None:
            prs.core_properties.author = author
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        return path

    return _make


# --------------------------------------------------------------------------
# pdf ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_pdf(tmp_path: Path):
    def _make(
        name: str = "doc.pdf",
        *,
        pages: list[dict] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        doc = fitz.open()
        for spec in pages or [{}]:
            pg = doc.new_page()
            for text, pos in spec.get("texts", []):
                pg.insert_text(pos, text, fontsize=11)
            # 罫線付きの表を描く: grid = 行数・列数・原点・セルサイズ・セルテキスト
            grid = spec.get("grid")
            if grid:
                n_rows = grid["rows"]
                n_cols = grid["cols"]
                x0, y0 = grid.get("origin", (100, 200))
                cw = grid.get("cw", 80)
                ch = grid.get("ch", 30)
                for i in range(n_rows + 1):
                    pg.draw_line((x0, y0 + i * ch), (x0 + n_cols * cw, y0 + i * ch))
                for j in range(n_cols + 1):
                    pg.draw_line((x0 + j * cw, y0), (x0 + j * cw, y0 + n_rows * ch))
                for (i, j), val in grid.get("cells", {}).items():
                    pg.insert_text((x0 + j * cw + 4, y0 + i * ch + 18), val, fontsize=10)
            for img, rect in spec.get("images", []):
                pg.insert_image(fitz.Rect(*rect), stream=img)
        if title is not None:
            doc.set_metadata({**(doc.metadata or {}), "title": title})
        if author is not None:
            doc.set_metadata({**(doc.metadata or {}), "author": author})
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(path))
        doc.close()
        return path

    return _make
