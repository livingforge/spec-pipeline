"""PowerPoint (.pptx) の抽出器。

スライドごとにシェイプを走査し、テキストフレーム・表・画像を抽出する。
グループシェイプは再帰的に展開する。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..models import ExtractionResult, ImageElement, TableElement, TextElement
from .base import ImageSaver


def extract_pptx(path: Path, saver: ImageSaver) -> ExtractionResult:
    prs = Presentation(str(path))
    result = ExtractionResult(source=path.name, file_type="pptx")

    props = prs.core_properties
    result.metadata = {
        "title": props.title or None,
        "author": props.author or None,
        "created": props.created.isoformat() if props.created else None,
        "modified": props.modified.isoformat() if props.modified else None,
        "slide_count": len(prs.slides),
    }

    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            _extract_shape(shape, i, saver, result)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                result.elements.append(
                    TextElement(content=notes, style="notes", location={"slide": i})
                )
    return result


def _extract_shape(shape, slide_no: int, saver: ImageSaver, result: ExtractionResult) -> None:
    loc = {"slide": slide_no, "shape_name": shape.name}

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _extract_shape(child, slide_no, saver, result)
        return

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        rel_path = saver.save(image.blob, image.ext)
        width, height = image.size
        result.elements.append(
            ImageElement(
                file=rel_path, format=image.ext, width=width, height=height, location=loc
            )
        )
        return

    if shape.has_table:
        rows = [
            [cell.text.strip() for cell in row.cells]
            for row in shape.table.rows
        ]
        result.elements.append(TableElement(rows=rows, location=loc))
        return

    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            result.elements.append(TextElement(content=text, location=loc))
